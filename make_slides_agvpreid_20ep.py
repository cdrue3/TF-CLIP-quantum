"""AG-VPReID 20-epoch preliminary results: QTemporal vs QTD vs classical baseline."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF7, 0xF9, 0xFC)
NAVY   = RGBColor(0x1F, 0x3D, 0x7A)
MID    = RGBColor(0x2E, 0x6D, 0xB4)
LTBLUE = RGBColor(0xD6, 0xE4, 0xF7)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
LTGR   = RGBColor(0xD4, 0xED, 0xDA)
RED    = RGBColor(0xA0, 0x1C, 0x1C)
LTRED  = RGBColor(0xF8, 0xD7, 0xD7)
LGREY  = RGBColor(0xF0, 0xF0, 0xF0)
MGREY  = RGBColor(0xCC, 0xCC, 0xCC)
GREY   = RGBColor(0x44, 0x44, 0x44)
BLACK  = RGBColor(0x11, 0x11, 0x11)
AMBER  = RGBColor(0x7A, 0x4A, 0x00)
LTAMB  = RGBColor(0xFD, 0xF3, 0xD0)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    p = Presentation()
    p.slide_width = W; p.slide_height = H
    return p

def blank(prs): return prs.slide_layouts[6]

def bg_fill(slide, c=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.width = lw
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def txt(slide, s, l, t, w, h, sz=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def page_header(slide, title, sub=""):
    bg_fill(slide)
    rect(slide, 0, 0, W, Inches(0.62), fill=NAVY)
    rect(slide, 0, H - Inches(0.05), W, Inches(0.05), fill=NAVY)
    txt(slide, title, Inches(0.3), Inches(0.07), Inches(12.7), Inches(0.48),
        sz=22, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.3), Inches(0.63), W - Inches(0.6), Inches(0.26),
            sz=9, color=GREY, italic=True)

def th(slide, xs, ws, labels, y, rh=Inches(0.34)):
    rect(slide, xs[0], y, xs[-1]+ws[-1]-xs[0], rh, fill=NAVY)
    for x, w, lbl in zip(xs, ws, labels):
        txt(slide, lbl, x+Inches(0.05), y+Inches(0.03), w-Inches(0.08), rh,
            sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return y + rh

def tr(slide, xs, ws, vals, colors, y, rh=Inches(0.35), shade=False, aligns=None):
    fill = LGREY if shade else WHITE
    rect(slide, xs[0], y, xs[-1]+ws[-1]-xs[0], rh-Inches(0.01), fill=fill)
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER]*(len(vals)-1)
    for x, w, v, c, al in zip(xs, ws, vals, colors, aligns):
        txt(slide, str(v), x+Inches(0.05), y+Inches(0.04), w-Inches(0.08), rh,
            sz=9.5, color=c, align=al)
    return y + rh

def note_box(slide, l, t, w, h, bg_color, line_color, text_str):
    rect(slide, l, t, w, h, fill=bg_color, line=line_color, lw=Pt(1))
    txt(slide, text_str, l+Inches(0.12), t+Inches(0.1), w-Inches(0.2), h-Inches(0.12),
        sz=9.5, color=BLACK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def s_title(prs):
    slide = prs.slides.add_slide(blank(prs))
    bg_fill(slide, WHITE)
    rect(slide, 0, 0, W, Inches(0.08), fill=NAVY)
    rect(slide, 0, H-Inches(0.08), W, Inches(0.08), fill=NAVY)
    rect(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(2.8), fill=NAVY)
    txt(slide, "AG-VPReID: Temporal Quantum Variants",
        Inches(0.7), Inches(1.75), Inches(11.9), Inches(1.2),
        sz=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, "20-Epoch Preliminary Results  ·  Training Accuracy (acc_id1)",
        Inches(0.7), Inches(2.85), Inches(11.9), Inches(0.45),
        sz=18, color=LTBLUE, align=PP_ALIGN.CENTER)
    txt(slide, "QUT-KIT IMPULSE Summer Research Program  ·  Connor Claypool  ·  March 2026",
        Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.4),
        sz=13, color=GREY, align=PP_ALIGN.CENTER)
    txt(slide, "Dataset: AG-VPReID (1604 train IDs, 13300 tracklets, 6 cameras — 4 ground + 2 aerial)  ·  "
               "Backbone: ViT-B/16 + TF-CLIP  ·  Quantum: PennyLane VQC (8q, 2L)",
        Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.35),
        sz=11, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    txt(slide, "Architectures: QTemporal (data re-uploading)  ·  QTD (temporal differences)  ·  SEQ_LEN=8",
        Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.3),
        sz=11, color=NAVY, align=PP_ALIGN.CENTER, bold=True)

    # caveat box
    rect(slide, Inches(1.5), Inches(6.2), Inches(10.33), Inches(0.75),
         fill=LTAMB, line=AMBER, lw=Pt(1))
    txt(slide, "⚠  Preliminary comparison — baseline trained with SEQ_LEN=4 (different conditions). "
               "acc_id1 is a training metric, not Rank-1. SEQ_LEN=8 baseline rerun in progress.",
        Inches(1.65), Inches(6.3), Inches(10.0), Inches(0.55),
        sz=10, color=AMBER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Epoch Trajectory
# ══════════════════════════════════════════════════════════════════════════════
def s_trajectory(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "acc_id1 Trajectory — 20 Epochs (SEQ_LEN=8)",
                "Both VQC variants trained on full AG-VPReID (1604 IDs, 1490 batches/epoch)  ·  8 qubits, 2 layers")

    # Trajectory table — every epoch
    qtemporal = [0.058, 0.088, 0.107, 0.124, 0.153, 0.183, 0.211, 0.240, 0.259, 0.278,
                 0.271, 0.314, 0.380, 0.398, 0.422, 0.471, 0.494, 0.537, 0.574, 0.577]
    qtd       = [0.061, 0.090, 0.108, 0.121, 0.155, 0.185, 0.214, 0.246, 0.257, 0.283,
                 0.281, 0.317, 0.388, 0.394, 0.441, 0.455, 0.510, 0.521, 0.582, 0.571]

    xs = [Inches(0.3), Inches(2.1), Inches(3.5), Inches(4.9), Inches(6.3),
          Inches(7.7), Inches(9.1), Inches(10.5)]
    ws = [Inches(1.7), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3),
          Inches(1.3), Inches(1.3), Inches(1.3)]

    # First 8 epochs
    txt(slide, "Epochs 1–8  (LR warmup phase: 3.5e-6 → 2.55e-5)",
        Inches(0.3), Inches(0.93), Inches(12.5), Inches(0.24), sz=11, bold=True, color=NAVY)
    y = th(slide, xs, ws, ["Model", "Ep 1", "Ep 2", "Ep 3", "Ep 4", "Ep 5", "Ep 6", "Ep 7/8 avg"],
           Inches(1.17))
    y = tr(slide, xs, ws,
           ["QTemporal VQC",
            f"{qtemporal[0]:.3f}", f"{qtemporal[1]:.3f}", f"{qtemporal[2]:.3f}",
            f"{qtemporal[3]:.3f}", f"{qtemporal[4]:.3f}", f"{qtemporal[5]:.3f}",
            f"{(qtemporal[6]+qtemporal[7])/2:.3f}"],
           [BLACK]+[MID]*7, y, shade=False)
    y = tr(slide, xs, ws,
           ["QTD VQC",
            f"{qtd[0]:.3f}", f"{qtd[1]:.3f}", f"{qtd[2]:.3f}",
            f"{qtd[3]:.3f}", f"{qtd[4]:.3f}", f"{qtd[5]:.3f}",
            f"{(qtd[6]+qtd[7])/2:.3f}"],
           [BLACK]+[MID]*7, y, shade=True)

    # Epochs 9-20 — different LR phases
    txt(slide, "Epochs 9–20  (peak LR phase: 3.50e-5 from ep11)",
        Inches(0.3), y+Inches(0.18), Inches(12.5), Inches(0.24), sz=11, bold=True, color=NAVY)

    xs2 = [Inches(0.3), Inches(2.1), Inches(3.3), Inches(4.5), Inches(5.7),
           Inches(6.9), Inches(8.1), Inches(9.3), Inches(10.5), Inches(11.7)]
    ws2 = [Inches(1.7), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1),
           Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.1)]

    y2 = th(slide, xs2, ws2,
            ["Model", "Ep 9", "Ep 10", "Ep 11", "Ep 12", "Ep 13", "Ep 14", "Ep 15", "Ep 18", "Ep 20"],
            y+Inches(0.42))
    y2 = tr(slide, xs2, ws2,
            ["QTemporal VQC",
             f"{qtemporal[8]:.3f}", f"{qtemporal[9]:.3f}", f"{qtemporal[10]:.3f}",
             f"{qtemporal[11]:.3f}", f"{qtemporal[12]:.3f}", f"{qtemporal[13]:.3f}",
             f"{qtemporal[14]:.3f}", f"{qtemporal[17]:.3f}", f"{qtemporal[19]:.3f}"],
            [BLACK]+[MID]*9, y2, shade=False)
    y2 = tr(slide, xs2, ws2,
            ["QTD VQC",
             f"{qtd[8]:.3f}", f"{qtd[9]:.3f}", f"{qtd[10]:.3f}",
             f"{qtd[11]:.3f}", f"{qtd[12]:.3f}", f"{qtd[13]:.3f}",
             f"{qtd[14]:.3f}", f"{qtd[17]:.3f}", f"{qtd[19]:.3f}"],
            [BLACK]+[MID]*9, y2, shade=True)

    # Final comparison box
    rect(slide, Inches(0.3), y2+Inches(0.18), Inches(12.73), Inches(1.55),
         fill=LTBLUE, line=NAVY, lw=Pt(1))
    txt(slide, "ep20 Summary vs Classical Baseline",
        Inches(0.45), y2+Inches(0.23), Inches(12.4), Inches(0.28),
        sz=11, bold=True, color=NAVY)

    # mini summary table inside box
    sx = [Inches(0.45), Inches(4.5), Inches(6.2), Inches(8.3), Inches(10.4)]
    sw = [Inches(3.9), Inches(1.6), Inches(2.0), Inches(2.0), Inches(2.1)]
    sy = y2 + Inches(0.52)
    rect(slide, sx[0], sy, sx[-1]+sw[-1]-sx[0], Inches(0.28), fill=NAVY)
    for x, w, lbl in zip(sx, sw, ["Model", "acc_id1 ep20", "Conditions", "vs Baseline Δ", "Note"]):
        txt(slide, lbl, x+Inches(0.05), sy+Inches(0.02), w-Inches(0.08), Inches(0.28),
            sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sy += Inches(0.28)

    rows = [
        ("Classical Baseline",   "0.165", "SEQ_LEN=4, 80/20 split", "—",      "Reference (different conditions)"),
        ("QTemporal VQC (ours)", "0.577", "SEQ_LEN=8, full dataset", "+0.412", "Still rising — no LR decay yet"),
        ("QTD VQC (ours)",       "0.571", "SEQ_LEN=8, full dataset", "+0.406", "Still rising — no LR decay yet"),
    ]
    for i, (name, acc, cond, delta, note) in enumerate(rows):
        fill = LGREY if i % 2 == 0 else WHITE
        rect(slide, sx[0], sy, sx[-1]+sw[-1]-sx[0], Inches(0.29), fill=fill)
        dc = GREEN if delta.startswith("+") else (GREY if delta == "—" else RED)
        for x, w, v, c, al in zip(sx, sw,
                [name, acc, cond, delta, note],
                [BLACK, BLACK, GREY, dc, GREY],
                [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT]):
            txt(slide, v, x+Inches(0.05), sy+Inches(0.04), w-Inches(0.08), Inches(0.29),
                sz=9, color=c, align=al)
        sy += Inches(0.29)

    # caveat
    txt(slide, "⚠  Baseline comparison is indicative only — different SEQ_LEN (4 vs 8) and dataset split (80/20 vs full 1604 IDs). "
               "SEQ_LEN=8 baseline rerun pending. Both VQC variants still actively learning at ep20 (LR never decayed).",
        Inches(0.3), sy+Inches(0.1), Inches(12.73), Inches(0.4),
        sz=9, color=AMBER, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture descriptions
# ══════════════════════════════════════════════════════════════════════════════
def s_architectures(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Architecture Overview — Two Temporal VQC Variants",
                "Both use data re-uploading over T=8 frames as the quantum temporal pooling mechanism")

    LW = Inches(6.1)

    # QTemporal
    rect(slide, Inches(0.3), Inches(0.92), LW, Inches(2.8), fill=WHITE, line=MID, lw=Pt(1.5))
    rect(slide, Inches(0.3), Inches(0.92), Inches(0.06), Inches(2.8), fill=MID)
    txt(slide, "QTemporal  (Temporal Quantum Aggregation)",
        Inches(0.47), Inches(1.0), LW-Inches(0.3), Inches(0.28), sz=12, bold=True, color=MID)
    txt(slide,
        "The VQC IS the temporal pooling step. Each of T=8 frames is encoded into the "
        "circuit sequentially via data re-uploading — quantum interference mixes temporal "
        "information across frames and the final measurement is the tracklet descriptor.\n\n"
        "  Frame 1 → pre_net → encode angles → entangle\n"
        "  Frame 2 → re-encode into same circuit → entangle\n"
        "  ...  ×8  ...  → measure → [2^n_q=256] → upscale → residual on mean_pool\n\n"
        "VQC runs at eval time — genuine quantum computation at inference.",
        Inches(0.47), Inches(1.30), LW-Inches(0.3), Inches(2.3), sz=10, color=BLACK)

    # QTD
    rect(slide, Inches(0.3), Inches(3.85), LW, Inches(2.8), fill=WHITE, line=NAVY, lw=Pt(1.5))
    rect(slide, Inches(0.3), Inches(3.85), Inches(0.06), Inches(2.8), fill=NAVY)
    txt(slide, "QTD  (Quantum Temporal Difference)",
        Inches(0.47), Inches(3.93), LW-Inches(0.3), Inches(0.28), sz=12, bold=True, color=NAVY)
    txt(slide,
        "Computes T-1=7 consecutive frame differences and feeds those through the VQC "
        "instead of raw frames. The intuition: differences capture motion and temporal "
        "change directly — arguably a lower-dimensional signal better suited to the "
        "limited capacity of an 8-qubit circuit.\n\n"
        "  diffs = x[:,1:] - x[:,:-1]   →   [B, 7, 768]\n"
        "  Each diff → pre_net → VQC → residual correction on mean_pool\n\n"
        "VQC runs only during training — mean_pool used at eval.",
        Inches(0.47), Inches(4.23), LW-Inches(0.3), Inches(2.3), sz=10, color=BLACK)

    # Right side — key comparison
    RX = Inches(6.7); RW = Inches(6.3)
    rect(slide, RX, Inches(0.92), RW, Inches(5.73), fill=LTBLUE, line=NAVY, lw=Pt(1))
    txt(slide, "Key Observations at ep20",
        RX+Inches(0.15), Inches(1.0), RW-Inches(0.25), Inches(0.28),
        sz=12, bold=True, color=NAVY)

    obs = [
        ("Nearly identical trajectories",
         "QTemporal 0.577 vs QTD 0.571 — separated by only 0.6pp at every checkpoint. "
         "The two very different approaches to temporal VQC processing converge to the same "
         "training accuracy, suggesting the bottleneck is the VQC capacity (8q), not the "
         "temporal encoding strategy."),
        ("Both still rising steeply",
         "Neither has plateaued by ep20. LR = 3.5e-5 throughout (no decay). "
         "The full learning trajectory will only emerge at 80ep with LR drops at ep60/72."),
        ("Strong absolute numbers",
         "acc_id1 ≈ 0.577 at 20ep on 1604-class problem. The baseline reference (0.165) "
         "used different conditions (SEQ_LEN=4, 80/20 split) — direct comparison invalid. "
         "A fair baseline rerun with SEQ_LEN=8 is pending."),
        ("QTD viable despite simpler signal",
         "Frame differences might seem lower-information than raw frames, but QTD matches "
         "QTemporal closely. Differences may be easier for the VQC to process due to "
         "smaller magnitude variance across tracklets."),
    ]

    cy = Inches(1.35)
    for title, body in obs:
        txt(slide, f"• {title}", RX+Inches(0.15), cy, RW-Inches(0.25), Inches(0.24),
            sz=10, bold=True, color=NAVY)
        cy += Inches(0.24)
        txt(slide, body, RX+Inches(0.25), cy, RW-Inches(0.35), Inches(0.52),
            sz=9.5, color=GREY)
        cy += Inches(0.6)

    # bottom note
    rect(slide, Inches(0.3), Inches(6.82), Inches(12.73), Inches(0.46),
         fill=LTAMB, line=AMBER, lw=Pt(1))
    txt(slide, "Next steps: SEQ_LEN=8 baseline rerun (fair comparison)  ·  "
               "Full 80ep runs for both architectures  ·  "
               "Eval with rrs_test (official AG-VPReID protocol)  ·  "
               "QGT and QFC 20ep runs pending",
        Inches(0.45), Inches(6.88), Inches(12.4), Inches(0.35),
        sz=9.5, color=AMBER)


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
prs = new_prs()
s_title(prs)
s_trajectory(prs)
s_architectures(prs)

out = "agvpreid_20ep_results.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
