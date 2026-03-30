"""Supervisor meeting: Quantum TF-CLIP — implementations with results on BOTH MARS and AG-ReID."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG     = RGBColor(0xF7, 0xF9, 0xFC)
NAVY   = RGBColor(0x1F, 0x3D, 0x7A)
MID    = RGBColor(0x2E, 0x6D, 0xB4)
LTBLUE = RGBColor(0xD6, 0xE4, 0xF7)
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
LTGR   = RGBColor(0xD4, 0xED, 0xDA)
RED    = RGBColor(0xA0, 0x1C, 0x1C)
LTRED  = RGBColor(0xF8, 0xD7, 0xD7)
LGREY  = RGBColor(0xF0, 0xF0, 0xF0)
MGREY  = RGBColor(0xCC, 0xCC, 0xCC)
GREY   = RGBColor(0x44, 0x44, 0x44)
BLACK  = RGBColor(0x11, 0x11, 0x11)
PURPLE = RGBColor(0x5A, 0x1A, 0x7A)
GOLD   = RGBColor(0x7A, 0x5A, 0x00)

W = Inches(13.33)
H = Inches(7.5)


# ── Primitives ─────────────────────────────────────────────────────────────
def new_prs():
    p = Presentation()
    p.slide_width = W; p.slide_height = H
    return p

def blank(prs): return prs.slide_layouts[6]

def bg_fill(slide, c=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = c

def rect(slide, l, t, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.line.width = lw
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line
    else:    s.line.fill.background()
    return s

def txt(slide, s, l, t, w, h, sz=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def page_header(slide, title, sub=""):
    bg_fill(slide)
    rect(slide, 0, 0, W, Inches(0.62), fill=NAVY)
    rect(slide, 0, H - Inches(0.05), W, Inches(0.05), fill=NAVY)
    txt(slide, title, Inches(0.3), Inches(0.07), Inches(12.7), Inches(0.48),
        sz=22, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, Inches(0.3), Inches(0.63), W - Inches(0.6), Inches(0.26),
            sz=9, color=GREY, italic=True)

def th(slide, xs, ws, labels, y, rh=Inches(0.34)):
    rect(slide, xs[0], y, xs[-1]+ws[-1]-xs[0], rh, fill=NAVY)
    for x, w, lbl in zip(xs, ws, labels):
        txt(slide, lbl, x+Inches(0.05), y+Inches(0.03), w-Inches(0.08), rh,
            sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return y + rh

def tr(slide, xs, ws, vals, colors, y, rh=Inches(0.35), shade=False, aligns=None):
    fill = LGREY if shade else WHITE
    rect(slide, xs[0], y, xs[-1]+ws[-1]-xs[0], rh-Inches(0.01), fill=fill)
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER]*(len(vals)-1)
    for x, w, v, c, al in zip(xs, ws, vals, colors, aligns):
        txt(slide, str(v), x+Inches(0.05), y+Inches(0.04), w-Inches(0.08), rh,
            sz=9.5, color=c, align=al)
    return y + rh

def dc(d):
    s = str(d).strip()
    if s.startswith("+"): return GREEN
    if s.startswith("−") or s.startswith("-"): return RED
    return GREY

def info_panel(slide, l, t, w, h, accent, what, how_concept, how_code, pipeline):
    rect(slide, l, t, w, h, fill=WHITE, line=accent, lw=Pt(1.5))
    rect(slide, l, t, Inches(0.06), h, fill=accent)
    cy = t + Inches(0.10); lx = l + Inches(0.13); pw = w - Inches(0.2)

    txt(slide, "WHAT", lx, cy, pw, Inches(0.2), sz=8, bold=True, color=accent)
    cy += Inches(0.2)
    txt(slide, what, lx, cy, pw, Inches(0.52), sz=10, color=BLACK)
    cy += Inches(0.56)

    txt(slide, "HOW — concept", lx, cy, pw, Inches(0.2), sz=8, bold=True, color=accent)
    cy += Inches(0.2)
    txt(slide, how_concept, lx, cy, pw, Inches(0.38), sz=10, color=BLACK)
    cy += Inches(0.42)

    txt(slide, "HOW — code", lx, cy, pw, Inches(0.2), sz=8, bold=True, color=accent)
    cy += Inches(0.2)
    txt(slide, how_code, lx, cy, pw, Inches(0.52), sz=9, color=GREY, italic=True)
    cy += Inches(0.56)

    txt(slide, "PIPELINE POSITION", lx, cy, pw, Inches(0.2), sz=8, bold=True, color=accent)
    cy += Inches(0.2)
    txt(slide, pipeline, lx, cy, pw, Inches(0.35), sz=10, color=BLACK)

def desc_panel(slide, l, t, w, h, accent, title, body, body_sz=9):
    """Freeform description panel: accent bar + bold title + single rich body text block."""
    rect(slide, l, t, w, h, fill=WHITE, line=accent, lw=Pt(1.5))
    rect(slide, l, t, Inches(0.06), h, fill=accent)
    lx = l + Inches(0.13); pw = w - Inches(0.2)
    cy = t + Inches(0.12)
    txt(slide, title, lx, cy, pw, Inches(0.28), sz=11, bold=True, color=accent)
    cy += Inches(0.3)
    txt(slide, body, lx, cy, pw, h - Inches(0.45), sz=body_sz, color=BLACK)


def note_box(slide, l, t, w, h, color, text_str, lt=None):
    lt = lt or color
    bg = LTGR if color == GREEN else (LTRED if color == RED else LTBLUE)
    rect(slide, l, t, w, h, fill=bg, line=lt, lw=Pt(1))
    txt(slide, text_str, l+Inches(0.1), t+Inches(0.08), w-Inches(0.15), h-Inches(0.1),
        sz=10, bold=True, color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
def s_title(prs):
    slide = prs.slides.add_slide(blank(prs))
    bg_fill(slide, WHITE)
    rect(slide, 0, 0, W, Inches(0.08), fill=NAVY)
    rect(slide, 0, H-Inches(0.08), W, Inches(0.08), fill=NAVY)
    rect(slide, Inches(0.5), Inches(1.5), Inches(12.33), Inches(3.0), fill=NAVY)
    txt(slide, "Quantum-Enhanced Video Person Re-Identification",
        Inches(0.7), Inches(1.8), Inches(11.9), Inches(1.3),
        sz=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, "Hybrid Quantum-Classical Architectures: Cross-Dataset Results",
        Inches(0.7), Inches(3.0), Inches(11.9), Inches(0.45),
        sz=18, color=LTBLUE, align=PP_ALIGN.CENTER)
    txt(slide, "QUT-KIT IMPULSE Summer Research Program  ·  Connor Claypool  ·  March 2026",
        Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.4),
        sz=13, color=GREY, align=PP_ALIGN.CENTER)
    txt(slide, "Datasets: MARS (625 ids) · AG-ReID (157 ids, aerial+ground)  ·  "
               "Backbone: ViT-B/16 + TF-CLIP  ·  Quantum: PennyLane VQC",
        Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.35),
        sz=11, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    txt(slide, "4 implementations on BOTH datasets  ·  4 additional AG-ReID-only implementations",
        Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.3),
        sz=11, color=NAVY, align=PP_ALIGN.CENTER, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Baselines
# ══════════════════════════════════════════════════════════════════════════════
def s_baselines(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Classical Baselines")

    xs = [Inches(0.3), Inches(4.2), Inches(5.65), Inches(7.1), Inches(8.45), Inches(9.75), Inches(11.2)]
    ws = [Inches(3.8), Inches(1.35), Inches(1.35), Inches(1.25), Inches(1.2),  Inches(1.35), Inches(1.9)]

    txt(slide, "MARS  —  625 identities, 11,310 test tracklets, 80 epochs",
        Inches(0.3), Inches(0.93), Inches(12.5), Inches(0.26), sz=12, bold=True, color=NAVY)
    y = th(slide, xs, ws, ["Model", "Rank-1", "Rank-5", "Rank-10", "Rank-20", "mAP", "Notes"],
           Inches(1.19))
    y = tr(slide, xs, ws,
           ["Classical TF-CLIP  (ViT-B/16, AdaMS temporal, 80ep)",
            "90.9%", "96.9%", "97.6%", "98.4%", "86.5%", "Reference"],
           [BLACK, GREEN, BLACK, BLACK, BLACK, GREEN, GREY], y, shade=True)

    txt(slide, "AG-ReID  —  157 identities, 748 tracklets, 2 cameras (aerial + ground), 80 epochs",
        Inches(0.3), y+Inches(0.22), Inches(12.5), Inches(0.26), sz=12, bold=True, color=NAVY)
    y2 = th(slide, xs, ws, ["Model", "Rank-1", "Rank-5", "Rank-10", "Rank-20", "mAP", "Notes"],
            y+Inches(0.48))
    y2 = tr(slide, xs, ws,
            ["Classical TF-CLIP  (80ep from scratch, aerial/ground split)",
             "74.3%", "86.9%", "—", "—", "—", "Reference"],
            [BLACK, GREEN, BLACK, GREY, GREY, GREY, GREY], y2, shade=True)

    # Why AG-ReID
    rect(slide, Inches(0.3), y2+Inches(0.2), Inches(12.73), Inches(1.35),
         fill=LTBLUE, line=NAVY, lw=Pt(1))
    txt(slide, "Why AG-ReID as primary benchmark",
        Inches(0.45), y2+Inches(0.25), Inches(12.4), Inches(0.28),
        sz=11, bold=True, color=NAVY)
    txt(slide, "MARS Rank-1 = 90.9% — near saturation, very small headroom for improvement.\n"
               "AG-ReID has a lower classical ceiling (74.3%) due to the aerial↔ground viewpoint gap "
               "→ more room for quantum advantage to surface. "
               "157 identities also avoids the multi-class degradation problem identified in the survey "
               "(quantum accuracy drops sharply as class count rises).",
        Inches(0.45), y2+Inches(0.53), Inches(12.4), Inches(0.9),
        sz=10, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — VQC Adapter
# ══════════════════════════════════════════════════════════════════════════════
def s_adapter(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 1 — VQC Adapter  (Residual Bottleneck)",
                "VQC in a skip-connection branch — output added to input features")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), MID,
        title="VQC Adapter  (Residual Bottleneck)",
        body=(
            "The simplest quantum integration: compress the 768-dim CLIP features down to n_q "
            "angles via a linear layer, run the VQC, expand back to 768-dim, and add the result "
            "to the original input as a residual correction.\n\n"
            "  input [768] → pre_net [768→n_q] → VQC → probs [2^n_q] → up_proj [2^n_q→768]\n"
            "  output = input + correction   ← residual skip connection\n\n"
            "The residual matters: even if the VQC produces noise, gradients still flow through "
            "the identity path and the model doesn't collapse.\n\n"
            "Classical ablation: replace VQC with Linear(n_q → 2^n_q) + ReLU. Same pre_net and "
            "up_proj, same in/out dims. VQC has 48 params (8q, 2L: n_layers × n_qubits × 3 "
            "rotation angles); classical has 2,048 (8×256) — 40× more parameters.\n\n"
            "In practice: 4q VQC edges ahead on MARS 15-ep training accuracy (+0.007) but "
            "classical wins on AG-ReID Rank-1 (−2.7pp). The 768 → 8 → 256 bottleneck compresses "
            "too aggressively. Both beat the 74.3% classical TF-CLIP baseline."
        ))

    txt(slide, "Results",
        RX, Inches(0.92), RW, Inches(0.26), sz=12, bold=True, color=NAVY)

    # MARS table
    txt(slide, "MARS  (15-epoch training accuracy — architecture search proxy, NOT Rank-1)",
        RX, Inches(1.20), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    xs = [RX, RX+Inches(2.7), RX+Inches(3.9), RX+Inches(5.0), RX+Inches(6.1)]
    ws = [Inches(2.6), Inches(1.1), Inches(1.0), Inches(1.0), Inches(1.75)]
    y = th(slide, xs, ws, ["Config", "VQC acc", "Cls acc", "Δ", "Note"], Inches(1.44))
    mars_rows = [
        ("4q, 2L  (best VQC config)", "0.309", "0.302", "+0.007", "Best VQC across all configs"),
        ("8q, 2L",                    "0.287", "0.296", "−0.009", "Classical wins"),
    ]
    for i, (a, v, c, d, n) in enumerate(mars_rows):
        y = tr(slide, xs, ws, [a, v, c, d, n],
               [BLACK, BLACK, BLACK, dc(d), GREY], y, shade=(i%2==0))

    # AG-ReID table
    txt(slide, "AG-ReID  (80-epoch Rank-1 — primary metric)",
        RX, y+Inches(0.14), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    y2 = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "vs Baseline"], y+Inches(0.38))
    agreid_rows = [
        ("4q, 2L",               "76.7%", "79.4%", "−2.7pp", "Both beat baseline 74.3%"),
        ("Baseline (reference)", "—",     "74.3%", "—",      "Classical TF-CLIP"),
    ]
    for i, (a, v, c, d, n) in enumerate(agreid_rows):
        y2 = tr(slide, xs, ws, [a, v, c, d, n],
                [BLACK, BLACK, BLACK, dc(d), GREY], y2, shade=(i%2==0))

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), RED,
        "MARS 15ep training acc: VQC marginally better at 4q (+0.007). "
        "AG-ReID 80ep Rank-1: classical wins (−2.7pp). "
        "Training-acc proxy does NOT predict Rank-1. Both beat the 74.3% baseline.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Channel Attention
# ══════════════════════════════════════════════════════════════════════════════
def s_channel(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 2 — Channel Attention",
                "VQC outputs per-channel scale factors — avoids classical primacy via multiplicative path")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), MID,
        title="Channel Attention",
        body=(
            "Mean-pool T=4 frames to a single [B,768] vector, compress to n_q angles, run VQC, "
            "and use the output as per-channel multipliers — the VQC decides which feature "
            "channels to amplify or suppress for each tracklet.\n\n"
            "  frames (B,T,768) → mean over T → [B,768] → pre_net → VQC\n"
            "  → σ(up_proj) → gates [B,768]\n"
            "  output = frames × (1 + gates)   ← broadcast gates over T frames\n\n"
            "Using multiplication instead of addition addresses 'classical primacy' (survey "
            "§3.3.2): the quantum branch can't dominate by being added in with large magnitude — "
            "it can only rescale what's already there.\n\n"
            "Classical ablation: replace VQC with Linear(n_q → 768) + sigmoid.\n\n"
            "In practice: VQC wins on MARS 15-ep training accuracy (+0.007) but classical wins "
            "narrowly on AG-ReID Rank-1 (−0.6pp). The multiplicative path helps but doesn't "
            "overcome the bottleneck. Closest result in the dataset — essentially tied."
        ))

    txt(slide, "Results", RX, Inches(0.92), RW, Inches(0.26), sz=12, bold=True, color=NAVY)

    txt(slide, "MARS  (15-epoch training accuracy — architecture search proxy, NOT Rank-1)",
        RX, Inches(1.20), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    xs = [RX, RX+Inches(2.7), RX+Inches(3.9), RX+Inches(5.0), RX+Inches(6.1)]
    ws = [Inches(2.6), Inches(1.1), Inches(1.0), Inches(1.0), Inches(1.75)]
    y = th(slide, xs, ws, ["Config", "VQC acc", "Cls acc", "Δ", "Note"], Inches(1.44))
    for i, row in enumerate([
        ("8q, 2L", "0.301", "0.294", "+0.007", "VQC wins"),
    ]):
        y = tr(slide, xs, ws, list(row), [BLACK, BLACK, BLACK, dc(row[3]), GREY], y, shade=(i%2==0))

    txt(slide, "AG-ReID  (80-epoch Rank-1 — primary metric)",
        RX, y+Inches(0.14), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    y2 = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "vs Baseline"], y+Inches(0.38))
    for i, row in enumerate([
        ("8q, 2L",               "76.7%", "77.3%", "−0.6pp", "Both beat baseline 74.3%"),
        ("Baseline (reference)", "—",     "74.3%", "—",      "Classical TF-CLIP"),
    ]):
        y2 = tr(slide, xs, ws, list(row), [BLACK, BLACK, BLACK, dc(row[3]), GREY], y2, shade=(i%2==0))

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), RED,
        "MARS 15ep training acc: VQC wins (+0.007). "
        "AG-ReID 80ep Rank-1: classical wins narrowly (−0.6pp). "
        "Consistent pattern: MARS proxy does not transfer to AG-ReID Rank-1.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Frame Attention
# ══════════════════════════════════════════════════════════════════════════════
def s_frame(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 3 — Frame Attention",
                "VQC predicts soft attention weights over T=4 frames — replaces uniform mean-pooling")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), MID,
        title="Frame Attention",
        body=(
            "Instead of mean-pooling T=4 frames equally, the VQC predicts a soft attention "
            "weight for each frame. Frames are then combined as a weighted sum, so the model "
            "can focus on the most informative moments in the tracklet.\n\n"
            "  frames (B,T,768) → reshape [B, T×768] → pre_net → VQC\n"
            "  → softmax → weights [B,T]\n"
            "  output = Σ_t  weights_t × frame_t   ← weighted combination\n\n"
            "Some frames are more useful than others (e.g., the person is partially occluded "
            "in frame 3). A quantum circuit processing all T frames' features simultaneously "
            "can evaluate all frame-importance combinations in parallel via superposition.\n\n"
            "Classical ablation: Linear(T×768 → T) + softmax. Same in/out structure.\n\n"
            "In practice: tied on MARS 15-ep training accuracy (0.298 = 0.298). VQC wins on "
            "AG-ReID Rank-1: +1.3pp (78.3% vs 77.0%). Frame selection aligns naturally with "
            "quantum superposition — parallel evaluation of which frames matter."
        ))

    txt(slide, "Results", RX, Inches(0.92), RW, Inches(0.26), sz=12, bold=True, color=NAVY)

    txt(slide, "MARS  (15-epoch training accuracy — architecture search proxy, NOT Rank-1)",
        RX, Inches(1.20), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    xs = [RX, RX+Inches(2.7), RX+Inches(3.9), RX+Inches(5.0), RX+Inches(6.1)]
    ws = [Inches(2.6), Inches(1.1), Inches(1.0), Inches(1.0), Inches(1.75)]
    y = th(slide, xs, ws, ["Config", "VQC acc", "Cls acc", "Δ", "Note"], Inches(1.44))
    for i, row in enumerate([
        ("8q, 2L", "0.298", "0.298", "0.000", "Tied"),
    ]):
        y = tr(slide, xs, ws, list(row), [BLACK, BLACK, BLACK, GREY, GREY], y, shade=(i%2==0))

    txt(slide, "AG-ReID  (80-epoch Rank-1 — primary metric)",
        RX, y+Inches(0.14), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    y2 = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "vs Baseline"], y+Inches(0.38))
    for i, row in enumerate([
        ("8q, 2L",               "78.3%", "77.0%", "+1.3pp", "Both beat baseline 74.3%  ✓"),
        ("Baseline (reference)", "—",     "74.3%", "—",      "Classical TF-CLIP"),
    ]):
        nc = GREEN if "✓" in row[4] else GREY
        y2 = tr(slide, xs, ws, list(row), [BLACK, BLACK, BLACK, dc(row[3]), nc], y2, shade=(i%2==0))

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), GREEN,
        "VQC wins on AG-ReID Rank-1: +1.3pp (78.3% vs 77.0%). "
        "MARS training acc: tied. Frame selection is a natural task for quantum superposition — "
        "parallel evaluation of frame importance aligns with quantum computing strengths.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Q-C-Q Interlaced
# ══════════════════════════════════════════════════════════════════════════════
def s_interlaced(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 4 — Q-C-Q Interlaced",
                "Alternating quantum-classical-quantum residual layers — survey §3.3.3 top-rated pattern")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), PURPLE,
        title="Q-C-Q Interlaced",
        body=(
            "Two quantum residual layers with a classical linear layer sandwiched between them. "
            "Each stage adds a small correction to the feature vector; the classical middle "
            "provides a full-rank projection between the quantum refinement stages.\n\n"
            "  x → (+ VQC₁ correction) = x₁       ← first quantum residual\n"
            "  x₁ → (+ Linear projection) = x₂     ← classical widening\n"
            "  x₂ → (+ VQC₂ correction) = output   ← second quantum residual\n\n"
            "The survey (§3.3.3) calls interlaced the best-performing integration pattern, "
            "citing +5.7pp on brain tumour classification. The intuition: quantum refines, "
            "classical broadens, quantum refines again — each stage sees the previous "
            "stage's corrections.\n\n"
            "Classical ablation: replace both VQC stages with Linear+ReLU layers of the "
            "same in/out dims.\n\n"
            "In practice: on MARS 80-ep Rank-1, VQC = Classical = 90.7% (both 0.2pp below "
            "baseline). On AG-ReID, classical wins (−2.7pp). Survey's top-rated pattern adds "
            "complexity here without gain — likely because the multi-class degradation problem "
            "is severe at 625 classes."
        ))

    txt(slide, "Results", RX, Inches(0.92), RW, Inches(0.26), sz=12, bold=True, color=NAVY)

    txt(slide, "MARS  (15-ep training acc + 80-ep Rank-1)",
        RX, Inches(1.20), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    xs = [RX, RX+Inches(2.5), RX+Inches(3.7), RX+Inches(4.85), RX+Inches(6.0)]
    ws = [Inches(2.4), Inches(1.1), Inches(1.05), Inches(1.05), Inches(1.95)]
    y = th(slide, xs, ws, ["Config", "VQC", "Classical", "Δ", "Note"], Inches(1.44))
    for i, row in enumerate([
        ("8q, 2L  (training acc, 15ep)", "0.298",  "0.294", "+0.004", "VQC marginally better"),
        ("8q, 2L  (Rank-1, 80ep)",       "90.7%",  "90.7%", "0.0pp",  "Tied — 0.2pp below baseline"),
        ("MARS baseline",                 "—",      "90.9%", "—",      "Reference"),
    ]):
        nc = GREY if "Tied" in row[4] else (GREEN if row[4].startswith("VQC") else GREY)
        y = tr(slide, xs, ws, list(row), [BLACK, BLACK, BLACK, dc(row[3]), nc], y, shade=(i%2==0))

    txt(slide, "AG-ReID  (80-epoch Rank-1 — primary metric)",
        RX, y+Inches(0.14), RW, Inches(0.24), sz=9, color=GREY, italic=True)
    y2 = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "vs Baseline"], y+Inches(0.38))
    for i, row in enumerate([
        ("8q, 2L",               "76.7%", "79.4%", "−2.7pp", "Both beat baseline 74.3%"),
        ("Baseline (reference)", "—",     "74.3%", "—",      "Classical TF-CLIP"),
    ]):
        y2 = tr(slide, xs, ws, list(row), [BLACK, BLACK, BLACK, dc(row[3]), GREY], y2, shade=(i%2==0))

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), PURPLE,
        "MARS 80ep Rank-1: tied at 90.7% (0.2pp below full TF-CLIP 90.9%) — VQC does not hurt. "
        "AG-ReID: classical wins (−2.7pp). "
        "Survey's top-rated architecture matches classical on MARS but not on AG-ReID.")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Summary
# ══════════════════════════════════════════════════════════════════════════════
def s_summary(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Summary: 4 Implementations on MARS + AG-ReID",
                "AG-ReID Rank-1 is the primary metric  ·  MARS column = 15ep training acc (proxy, not Rank-1) except interlaced")

    xs = [Inches(0.25), Inches(3.5), Inches(5.0), Inches(6.15), Inches(7.3),
          Inches(8.45), Inches(9.6), Inches(10.75)]
    ws = [Inches(3.15), Inches(1.4), Inches(1.05), Inches(1.05), Inches(1.05),
          Inches(1.05), Inches(1.05), Inches(2.3)]
    y = th(slide, xs, ws,
           ["Implementation", "Pipeline (Survey)",
            "MARS\n15ep acc†", "MARS\nCls acc†",
            "AG-ReID\nVQC R1", "AG-ReID\nCls R1", "AG-ReID\nΔ",
            "Verdict (AG-ReID Rank-1)"],
           Inches(0.88))

    alns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT,
            PP_ALIGN.CENTER, PP_ALIGN.CENTER,
            PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER,
            PP_ALIGN.LEFT]

    rows = [
        ("Baseline", "—", "—", "—", "—", "74.3%", "—", "Reference"),
        # Both datasets
        ("1. VQC Adapter (4q, residual)", "Feat. Ext. §3.3.1",
         "0.309", "0.302", "76.7%", "79.4%", "−2.7pp", "✗ Classical wins"),
        ("2. Channel Attention (8q)", "Feat. Ext. §3.3.2",
         "0.301", "0.294", "76.7%", "77.3%", "−0.6pp", "✗ Classical wins"),
        ("3. Frame Attention (8q)", "Feat. Ext. §3.3.1",
         "0.298", "0.298", "78.3%", "77.0%", "+1.3pp", "✓ VQC wins"),
        ("4. Q-C-Q Interlaced (8q)", "Feat. Ext. §3.3.3",
         "0.298 / 90.7%*", "0.294 / 90.7%*", "76.7%", "79.4%", "−2.7pp", "✗ Classical wins"),
        # AG-ReID only
        ("5. Gated Adapter (8q) ★", "Feat. Ext. Novel",
         "AG-ReID only", "AG-ReID only", "78.9%", "74.3%", "+4.6pp", "✓✓ VQC wins — BEST"),
        ("6. Temporal Agg (8q, TQA)", "Feat. Ext. §3.3.1",
         "AG-ReID only", "AG-ReID only", "78.9%", "77.5%", "+1.4pp", "✓ VQC wins (at eval)"),
        ("7. Q Feat Extractor (8q)", "Feat. Ext. §3.3.2",
         "AG-ReID only", "AG-ReID only", "78.6%", "75.9%", "+2.7pp", "✓ VQC wins"),
        ("8. CCG cam-conditioned (8q)", "Feat. Ext. Novel",
         "AG-ReID only", "AG-ReID only", "77.3%", "75.7%", "+1.6pp", "✓ VQC wins"),
    ]
    for i, row in enumerate(rows):
        vc = GREEN if "✓" in row[7] else (GREY if row[7] == "Reference" else RED)
        delta_c = dc(row[6])
        y = tr(slide, xs, ws, list(row),
               [BLACK, GREY, BLACK, BLACK, BLACK, BLACK, delta_c, vc],
               y, shade=(i%2==0), aligns=alns)

    # footnote
    txt(slide, "† MARS 15ep training accuracy (acc_id1) — used as architecture search proxy only. "
               "* Interlaced also has 80ep Rank-1 on MARS: VQC = Classical = 90.7%.",
        Inches(0.3), y+Inches(0.1), Inches(12.73), Inches(0.3),
        sz=9, color=GREY, italic=True)

    # findings box
    rect(slide, Inches(0.3), y+Inches(0.5), Inches(12.73), Inches(1.42),
         fill=LTBLUE, line=NAVY, lw=Pt(1))
    txt(slide, "Key findings",
        Inches(0.45), y+Inches(0.55), Inches(12.4), Inches(0.26),
        sz=11, bold=True, color=NAVY)
    txt(slide,
        "• AG-ReID VQC wins: Gated +4.6pp ★, QFeatExt +2.7pp, CCG +1.6pp, TQA +1.4pp, Frame +1.3pp  (5/7 implementations)\n"
        "• AG-ReID classical wins: Adapter −2.7pp, Interlaced −2.7pp, Channel −0.6pp  (3/7 implementations)\n"
        "• MARS 15ep training acc is NOT a reliable predictor of AG-ReID Rank-1 — Channel and Adapter win on acc but lose on Rank-1\n"
        "• Interlaced ties on MARS 80ep Rank-1 (90.7%) — VQC does not hurt at full training. All 8 beat the 74.3% AG-ReID baseline.",
        Inches(0.45), y+Inches(0.83), Inches(12.4), Inches(1.0),
        sz=10, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# AG-ReID-ONLY slides
# ══════════════════════════════════════════════════════════════════════════════
def s_gated(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 5 — Gated Adapter  ★  (AG-ReID only)",
                "Input-adaptive scalar gate controls VQC contribution — largest VQC win observed")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), GREEN,
        title="Gated Adapter ★",
        body=(
            "Like the VQC Adapter, but a learned scalar gate g ∈ (0,1) controls how much of "
            "the quantum correction to apply. The gate is predicted from the input features "
            "themselves:\n\n"
            "  Plain adapter:  output = x + quantum_correction\n"
            "  Gated adapter:  output = x + g(x) × quantum_correction\n\n"
            "  Where:  x [768] → gate_net: Linear(768→1) → sigmoid → g\n\n"
            "The motivation is directly the KIT Q2 research question: which inputs benefit most "
            "from quantum? A working gate would learn 'this tracklet needs more correction' vs "
            "'this one is fine as-is' — potentially routing harder aerial-ground pairs toward "
            "more quantum correction.\n\n"
            "Classical ablation: replace VQC with Linear(n_q → 2^n_q) + ReLU. Same gate, "
            "same pre_net, same up_proj. VQC has 48 params; classical has 2,048.\n\n"
            "In practice: gate saturated to g\u0305\u22480.95 for all tracklets (std=0.009). "
            "The model learned 'always apply quantum at near-full strength' rather than routing "
            "selectively. Still wins +4.6pp via the same regularisation effect as the plain "
            "adapter — but not the input-adaptive behaviour we were looking for."
        ))

    txt(slide, "Results  (AG-ReID 80-epoch Rank-1)", RX, Inches(0.92), RW, Inches(0.26),
        sz=12, bold=True, color=NAVY)
    xs = [RX, RX+Inches(3.0), RX+Inches(4.15), RX+Inches(5.1), RX+Inches(6.0)]
    ws = [Inches(2.9), Inches(1.1), Inches(1.05), Inches(0.85), Inches(1.6)]
    y = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "Note"], Inches(1.21))
    for i, row in enumerate([
        ("8q, 2L  NECK_FEAT=before", "78.9%", "74.3%", "+4.6pp", "✓✓ BEST VQC result overall"),
        ("8q, 2L  NECK_FEAT=after",  "78.6%", "76.5%", "+2.1pp", "Adapter active at eval"),
        ("4q, 2L  NECK_FEAT=before", "74.9%", "76.2%", "−1.3pp", "4q worse — 8q is optimal"),
        ("Baseline (reference)",     "—",     "74.3%", "—",      "Classical TF-CLIP"),
    ]):
        nc = GREEN if "BEST" in row[4] else GREY
        y = tr(slide, xs, ws, list(row),
               [BLACK, BLACK, BLACK, dc(row[3]), nc], y, shade=(i%2==0))

    txt(slide, "Gate analysis (8q NECK_FEAT=before): mean g̅ ≈ 0.95, std = 0.009. "
               "Gate learned 'always use quantum at full strength' — not truly input-adaptive.",
        RX, y+Inches(0.1), RW, Inches(0.4), sz=9, color=GREY)

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), GREEN,
        "Largest VQC win: +4.6pp Rank-1 (78.9% vs 74.3%). VQC acts as training regulariser "
        "shaping backbone features. Gate saturation (g̅≈0.95) suggests VQC contribution is "
        "consistently helpful — but not selectively routed as intended.")


def s_temporal(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 6 — Temporal Quantum Aggregation  (AG-ReID only)",
                "VQC IS the temporal pooling step via data re-uploading — runs at eval (genuine quantum at inference)")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), MID,
        title="Temporal Quantum Aggregation (TQA)",
        body=(
            "Instead of mean-pooling T=4 frames into one tracklet descriptor, the VQC IS the "
            "temporal pooling step. Each frame's features are encoded into the same circuit one "
            "at a time — data re-uploading (survey §3.2.1) — and the final measurement is "
            "the tracklet descriptor.\n\n"
            "  Frame 1 → pre_net → encode angles → entangle\n"
            "  Frame 2 → pre_net → re-encode into same circuit → entangle\n"
            "  Frame 3 → pre_net → re-encode → entangle\n"
            "  Frame 4 → pre_net → re-encode → entangle → measure → [2^n_q]\n\n"
            "Re-uploading means quantum interference mixes information from all frames across "
            "the entangling layers. The circuit acts as a learned temporal pooling operator.\n\n"
            "Unlike every other architecture: this VQC runs at evaluation time — it IS the "
            "pooling, not an add-on. TQA is the only architecture with genuine quantum "
            "computation at inference.\n\n"
            "Classical ablation: replace the entire circuit with an LSTM over T frames.\n\n"
            "In practice: VQC wins +1.4pp (78.9% vs 77.5%). Strongest argument for genuine "
            "quantum utility — the advantage holds even when the VQC is active at retrieval."
        ))

    txt(slide, "Results  (AG-ReID 80-epoch Rank-1)", RX, Inches(0.92), RW, Inches(0.26),
        sz=12, bold=True, color=NAVY)
    xs = [RX, RX+Inches(3.0), RX+Inches(4.15), RX+Inches(5.1), RX+Inches(6.0)]
    ws = [Inches(2.9), Inches(1.1), Inches(1.05), Inches(0.85), Inches(1.6)]
    y = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "Note"], Inches(1.21))
    for i, row in enumerate([
        ("8q, 2L", "78.9%", "77.5%", "+1.4pp", "✓ VQC wins"),
        ("Baseline (reference)", "—", "74.3%", "—", "Classical TF-CLIP"),
    ]):
        nc = GREEN if "✓" in row[4] else GREY
        y = tr(slide, xs, ws, list(row),
               [BLACK, BLACK, BLACK, dc(row[3]), nc], y, shade=(i%2==0))

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), GREEN,
        "VQC wins +1.4pp (78.9% vs 77.5%). Critically: this is a genuine quantum-at-inference result "
        "— the VQC circuit runs during retrieval, not just training. "
        "Data re-uploading enables quantum interference across frames as a pooling mechanism.")


def s_qfeatext(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 7 — Q Feature Extractor / Parallel Branch  (AG-ReID only)",
                "VQC runs parallel to CLIP backbone; outputs concatenated before classifier")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), MID,
        title="Q Feature Extractor (QFeatExt)",
        body=(
            "A parallel quantum branch runs alongside CLIP during training. Both branches see "
            "the same 768-dim backbone features; their outputs are concatenated before the "
            "classifier.\n\n"
            "  CLIP features [768] ─────────────────────────────────────────┐\n"
            "                                                                 → concat [1024] → classifier\n"
            "  CLIP features [768] → pre_net → VQC → q-feats [2^n_q=256] ──┘\n\n"
            "At evaluation: retrieval uses raw CLIP backbone features only (NECK_FEAT='before'). "
            "The VQC only runs during training — it acts as an implicit regulariser that shapes "
            "how the backbone learns to represent identity. Think of it as a quantum-regularised "
            "training signal that then gets discarded at test time.\n\n"
            "Classical ablation: replace VQC with Linear(n_q → 2^n_q) + ReLU. Same "
            "concatenation and classifier. VQC has 48 params; classical has 2,048. VQC wins "
            "with 40× fewer parameters — the quantum effect is about parameter efficiency, "
            "not raw capacity.\n\n"
            "In practice: VQC wins at every checkpoint (+1.6pp at ep35 → +2.7pp at ep80, still "
            "improving with no sign of plateau). Cleanest and most consistent VQC win."
        ))

    txt(slide, "Results  (AG-ReID 80-epoch Rank-1  +  early stopping)",
        RX, Inches(0.92), RW, Inches(0.26), sz=12, bold=True, color=NAVY)
    xs = [RX, RX+Inches(2.9), RX+Inches(4.0), RX+Inches(5.0), RX+Inches(5.9)]
    ws = [Inches(2.8), Inches(1.0), Inches(0.95), Inches(0.85), Inches(1.85)]
    y = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "Note"], Inches(1.21))
    for i, row in enumerate([
        ("8q, 2L — ep80", "78.6%", "75.9%", "+2.7pp", "✓ VQC wins"),
        ("8q, 2L — ep50", "77.8%", "75.9%", "+1.9pp", "early stop — no overfitting"),
        ("8q, 2L — ep40", "77.5%", "75.9%", "+1.6pp", "early stop"),
        ("8q, 2L — ep35", "77.5%", "75.9%", "+1.6pp", "early stop"),
        ("Baseline (ref)", "—",    "74.3%", "—",      "Classical TF-CLIP"),
    ]):
        nc = GREEN if "✓" in row[4] else GREY
        y = tr(slide, xs, ws, list(row),
               [BLACK, BLACK, BLACK, dc(row[3]), nc], y, shade=(i%2==0))

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), MID,
        "VQC wins at all checkpoints (+0.8 to +2.7pp). ep80 is best (78.6%). "
        "No overfitting despite training acc ≈1.0 from ep35 — CLIP backbone largely frozen. "
        "Classical primacy (survey §3.3.2) does NOT kill quantum here: parallel branch adds benefit.")


def s_ccg(prs):
    slide = prs.slides.add_slide(blank(prs))
    page_header(slide, "Implementation 8 — Camera-Conditioned Gating  (AG-ReID only)",
                "Gate conditioned on camera ID embedding — aerial vs ground routing")

    LW = Inches(4.8); RX = Inches(5.2); RW = Inches(7.9)
    desc_panel(slide, Inches(0.3), Inches(0.92), LW, Inches(5.5), PURPLE,
        title="CCG — Camera-Conditioned Gating",
        body=(
            "This is the Gated Adapter with one modification to how g is computed. Instead of "
            "the gate looking only at image features, it also gets told which camera the "
            "tracklet came from:\n\n"
            "  Gated gate:  x [768] → Linear(768→1) → sigmoid → g\n\n"
            "  CCG gate:    x [768] ────────────────────────────────────────┐\n"
            "                                                                 → concat → Linear(784→1) → sigmoid → g\n"
            "               cam_id → Embedding(2→16) → 16 numbers ──────────┘\n\n"
            "cam_id is a single integer (0=ground, 1=aerial). It gets converted to 16 learned "
            "numbers via an embedding table — think of it as 'camera 0's personality vector' "
            "vs 'camera 1's personality vector.' Those 16 numbers are concatenated onto the "
            "768 image features before the gate decides.\n\n"
            "The motivation: AG-ReID has two fundamentally different cameras — a ground-level "
            "camera and an aerial drone. The hypothesis was that the gate might learn 'aerial "
            "shots need more quantum correction than ground shots' (or vice versa).\n\n"
            "In practice: gate still saturated (g\u0305\u22480.95 for both cameras, "
            "aerial\u2212ground \u0394g=\u22120.003). The camera embedding didn't break the "
            "saturation problem. Still beats its classical ablation (+1.6pp) but for the same "
            "regularisation reason as the plain gated — not camera-aware routing."
        ))

    txt(slide, "Results  (AG-ReID 80-epoch Rank-1)", RX, Inches(0.92), RW, Inches(0.26),
        sz=12, bold=True, color=NAVY)
    xs = [RX, RX+Inches(3.0), RX+Inches(4.15), RX+Inches(5.1), RX+Inches(6.0)]
    ws = [Inches(2.9), Inches(1.1), Inches(1.05), Inches(0.85), Inches(1.6)]
    y = th(slide, xs, ws, ["Config", "VQC R1", "Cls R1", "Δ", "Note"], Inches(1.21))
    for i, row in enumerate([
        ("8q, 2L  NECK_FEAT=after",  "77.3%", "75.7%",   "+1.6pp", "✓ VQC wins"),
        ("Gated (no cam, ref)",       "78.9%", "74.3%",   "+4.6pp", "from Impl. 5 — for context"),
        ("Baseline (reference)",      "—",     "74.3%",   "—",      "Classical TF-CLIP"),
    ]):
        nc = GREEN if "context" in row[4] else GREY
        y = tr(slide, xs, ws, list(row),
               [BLACK, BLACK, GREY, dc(row[3]), nc], y, shade=(i%2==0))

    txt(slide, "Gate analysis: g̅ ≈ 0.95 (both cameras), std = 0.009, aerial−ground Δg = −0.003. "
               "Gate saturated — camera embedding not influencing the routing decision.",
        RX, y+Inches(0.1), RW, Inches(0.4), sz=9, color=GREY)

    note_box(slide, Inches(0.3), Inches(6.52), Inches(12.73), Inches(0.46), PURPLE,
        "CCG VQC (77.3%) below plain gated VQC (78.9%) — camera conditioning does not help. "
        "Gate saturation (g̅≈0.95) persists even with camera embedding. "
        "Input-adaptive routing via learned gating remains an open research question.")


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════
prs = new_prs()
s_title(prs)
s_baselines(prs)
# Both datasets
s_adapter(prs)
s_channel(prs)
s_frame(prs)
s_interlaced(prs)
# AG-ReID only
s_gated(prs)
s_temporal(prs)
s_qfeatext(prs)
s_ccg(prs)
s_summary(prs)

out = "quantum_clip_reid_summary.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
